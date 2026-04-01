"""
Company Intel - Clay.com SCOUT Research
2-Slide Executive Deck - Rebuilt with proper spacing and varied palette
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Distinct, high-contrast palette - no two adjacent elements share a color
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x00, 0x78, 0xD4)
TEAL = RGBColor(0x00, 0x8B, 0x8B)
PURPLE = RGBColor(0x68, 0x3A, 0xB7)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
FOREST = RGBColor(0x2E, 0x7D, 0x32)
CRIMSON = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DGRAY = RGBColor(0x33, 0x33, 0x33)
MGRAY = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xF2, 0xF2, 0xF2)
ZEBRA = RGBColor(0xF7, 0xF7, 0xF7)

SW = Inches(13.333)
SH = Inches(7.5)


def rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()


def card(s, l, t, w, h, accent=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = RGBColor(0xDE, 0xDE, 0xDE); sh.line.width = Pt(1)
    if accent:
        rect(s, l, t, w, Pt(5), accent)


def tx(s, l, t, w, h, text, sz=11, c=DGRAY, bold=False, al=PP_ALIGN.LEFT, fn="Segoe UI"):
    b = s.shapes.add_textbox(l, t, w, h)
    b.text_frame.word_wrap = True
    p = b.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = fn; p.alignment = al


def slide_1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = LGRAY

    # === HEADER ===
    rect(sl, Inches(0), Inches(0), SW, Inches(1.0), NAVY)
    rect(sl, Inches(0), Inches(0), SW, Pt(5), BLUE)
    tx(sl, Inches(0.8), Inches(0.12), Inches(4), Inches(0.5),
       "Clay", sz=34, c=WHITE, bold=True, fn="Segoe UI Semibold")
    tx(sl, Inches(0.8), Inches(0.58), Inches(6), Inches(0.3),
       "COMPANY OVERVIEW AND POSITIONING", sz=10, c=RGBColor(0x80, 0xD0, 0xD0), bold=True)
    tx(sl, Inches(8.5), Inches(0.18), Inches(4.5), Inches(0.3),
       "clay.com", sz=12, c=WHITE, al=PP_ALIGN.RIGHT)
    tx(sl, Inches(8.5), Inches(0.52), Inches(4.5), Inches(0.3),
       "SCOUT Framework  |  April 2026", sz=9, c=RGBColor(0x90, 0xA0, 0xB8), al=PP_ALIGN.RIGHT)

    # === METRICS ROW (5 cards, wider, spaced) ===
    metrics = [
        ("300K+", "GTM Teams", BLUE),
        ("150+", "Data Providers", TEAL),
        ("140M", "Monthly AI Runs", PURPLE),
        ("16", "Case Studies", ORANGE),
        ("26", "Customers Identified", FOREST),
    ]
    card_w = Inches(2.32)
    gap = Inches(0.15)
    start_x = Inches(0.5)
    for i, (num, label, accent) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        card(sl, x, Inches(1.2), card_w, Inches(0.85), accent)
        tx(sl, x, Inches(1.28), card_w, Inches(0.4),
           num, sz=22, c=accent, bold=True, al=PP_ALIGN.CENTER, fn="Segoe UI Semibold")
        tx(sl, x, Inches(1.65), card_w, Inches(0.3),
           label, sz=10, c=MGRAY, al=PP_ALIGN.CENTER)

    # === LEFT CARD: What They Do ===
    card(sl, Inches(0.5), Inches(2.25), Inches(5.9), Inches(2.7), BLUE)
    tx(sl, Inches(0.75), Inches(2.38), Inches(4), Inches(0.3),
       "WHAT THEY DO", sz=12, c=BLUE, bold=True)
    tx(sl, Inches(0.75), Inches(2.75), Inches(5.4), Inches(0.65),
       "AI-powered GTM data platform combining 150+ enrichment "
       "providers, intent signals, and CRM data. One person can run "
       "campaigns that previously required the entire GTM team.",
       sz=11, c=DGRAY)

    tx(sl, Inches(0.75), Inches(3.45), Inches(4), Inches(0.3),
       "KEY CAPABILITIES", sz=11, c=BLUE, bold=True)
    caps = [
        "Data enrichment from 150+ providers",
        "AI agent builder (Claygent) with MCP",
        "Real-time intent signal monitoring",
        "CRM sync and workflow orchestration",
        "Native email sequencer",
    ]
    for i, c in enumerate(caps):
        tx(sl, Inches(0.85), Inches(3.75 + i * 0.24), Inches(5.2), Inches(0.22),
           f"\u2022  {c}", sz=10, c=DGRAY)

    # === RIGHT CARD: Positioning + Buyers ===
    card(sl, Inches(6.6), Inches(2.25), Inches(6.3), Inches(2.7), PURPLE)
    tx(sl, Inches(6.85), Inches(2.38), Inches(5), Inches(0.3),
       "POSITIONING", sz=12, c=PURPLE, bold=True)
    tx(sl, Inches(6.85), Inches(2.75), Inches(5.8), Inches(0.55),
       '"Every GTM data point imaginable, in one place."\n'
       "Understand your total market, your customers, and the gap between them.",
       sz=11, c=DGRAY)

    tx(sl, Inches(6.85), Inches(3.38), Inches(5), Inches(0.3),
       "TARGET BUYERS", sz=11, c=PURPLE, bold=True)
    buyers = [
        ("Revenue Ops / Sales Ops", "Primary buyer"),
        ("Growth / Demand Gen", "Outbound + lead scoring"),
        ("Marketing Operations", "CRM hygiene + segmentation"),
        ("C-Suite (CEO, CRO)", "Executive champions"),
    ]
    for i, (role, ctx) in enumerate(buyers):
        y = Inches(3.7) + Inches(i * 0.28)
        tx(sl, Inches(6.95), y, Inches(2.5), Inches(0.25),
           f"\u2022  {role}", sz=10, c=DGRAY, bold=True)
        tx(sl, Inches(9.6), y, Inches(3.1), Inches(0.25),
           ctx, sz=10, c=MGRAY)

    # === BOTTOM BAR: Quote ===
    rect(sl, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.75), NAVY)
    tx(sl, Inches(0.8), Inches(5.2), Inches(9.5), Inches(0.35),
       '"We have 3x our enrichment rate with Clay. A game changer for marketing, data, and operations."',
       sz=12, c=WHITE)
    tx(sl, Inches(0.8), Inches(5.55), Inches(6), Inches(0.25),
       "Adam Wall, Head of Sales Ops, Anthropic", sz=10, c=RGBColor(0x80, 0xD0, 0xD0), bold=True)

    # === SEGMENTS ROW ===
    card(sl, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.6), None)
    tx(sl, Inches(0.75), Inches(6.12), Inches(2.2), Inches(0.25),
       "CUSTOMER SEGMENTS", sz=10, c=DGRAY, bold=True)
    segs = [
        ("AI / ML", PURPLE), ("SaaS", BLUE), ("Security", CRIMSON),
        ("Fintech", TEAL), ("Enterprise", FOREST),
    ]
    for i, (seg, color) in enumerate(segs):
        x = Inches(3.0) + Inches(i * 1.5)
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.14), Inches(1.3), Inches(0.3))
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
        tx(sl, x, Inches(6.15), Inches(1.3), Inches(0.28),
           seg, sz=9, c=WHITE, bold=True, al=PP_ALIGN.CENTER)

    # Footer
    tx(sl, Inches(0.8), Inches(6.85), Inches(12), Inches(0.2),
       "Source: clay.com  |  SCOUT Framework  |  ISV Intel", sz=8, c=MGRAY)


def slide_2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = LGRAY

    # === HEADER ===
    rect(sl, Inches(0), Inches(0), SW, Inches(1.0), NAVY)
    rect(sl, Inches(0), Inches(0), SW, Pt(5), ORANGE)
    tx(sl, Inches(0.8), Inches(0.12), Inches(8), Inches(0.5),
       "Clay  |  Customer Evidence Map", sz=30, c=WHITE, bold=True, fn="Segoe UI Semibold")
    tx(sl, Inches(0.8), Inches(0.58), Inches(10), Inches(0.3),
       "16 CASE STUDIES  WITH QUOTES   |   10 LOGO FEATURES   |   26 TOTAL CUSTOMERS",
       sz=10, c=RGBColor(0xFF, 0xCC, 0x80), bold=True)

    # === CASE STUDY TABLE ===
    card(sl, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5), ORANGE)
    tx(sl, Inches(0.75), Inches(1.3), Inches(5), Inches(0.3),
       "CASE STUDIES", sz=12, c=ORANGE, bold=True)

    # Column headers - left
    tx(sl, Inches(0.95), Inches(1.62), Inches(1.5), Inches(0.22),
       "CUSTOMER", sz=9, c=MGRAY, bold=True)
    tx(sl, Inches(2.7), Inches(1.62), Inches(1.8), Inches(0.22),
       "KEY RESULT", sz=9, c=MGRAY, bold=True)
    tx(sl, Inches(4.5), Inches(1.62), Inches(2.0), Inches(0.22),
       "CONTACT", sz=9, c=MGRAY, bold=True)
    # Column headers - right
    tx(sl, Inches(6.95), Inches(1.62), Inches(1.5), Inches(0.22),
       "CUSTOMER", sz=9, c=MGRAY, bold=True)
    tx(sl, Inches(8.7), Inches(1.62), Inches(1.8), Inches(0.22),
       "KEY RESULT", sz=9, c=MGRAY, bold=True)
    tx(sl, Inches(10.5), Inches(1.62), Inches(2.0), Inches(0.22),
       "CONTACT", sz=9, c=MGRAY, bold=True)

    # Divider line under headers
    rect(sl, Inches(0.7), Inches(1.85), Inches(5.55), Pt(1), RGBColor(0xCC, 0xCC, 0xCC))
    rect(sl, Inches(6.7), Inches(1.85), Inches(5.9), Pt(1), RGBColor(0xCC, 0xCC, 0xCC))

    left = [
        ("OpenAI", "2x enrichment", "Scotty Huhn, Rev. Strategy", BLUE),
        ("Anthropic", "3x enrichment rate", "Adam Wall, Head Sales Ops", PURPLE),
        ("Rippling", "2x email perf.", "Ryan Narod, VP Mktg", TEAL),
        ("Figma", "GTM orchestration", "Kyle Ketchum, Mktg Ops", ORANGE),
        ("Intercom", "Data consolidation", "A. DeMoulin, Dir. GTM Ops", FOREST),
        ("Verkada", "Auto enrichment", "D. Grieco, Dir. Growth", CRIMSON),
        ("Mistral AI", "25K accts / 2 wks", "Joel Davidson, VP RevOps", BLUE),
        ("Vanta", "GTM stack pillar", "Stevie Case, CRO", PURPLE),
    ]
    right = [
        ("ElevenLabs", "Auto lead scoring", "Raman Khanna, Growth", TEAL),
        ("Sendoso", "10x productivity", "Kris Rudegraap, CEO", ORANGE),
        ("Hex", "+50% win-rate", "B. Clancy, GTM Eng.", FOREST),
        ("Sana", "60% CRM accuracy", "F. Hillstrom, GTM Ops", CRIMSON),
        ("Rootly", "100% automation", "JJ Tang, CEO", BLUE),
        ("Legora", "+70% velocity", "F. Pavan, Dir. Demand Gen", PURPLE),
        ("Exit Five", "27K enriched", "Dan Murphy, COO", TEAL),
        ("depthfirst", "95%+ fill rates", "Mark Hardy, Head RevOps", ORANGE),
    ]

    row_h = Inches(0.35)
    row_gap = Inches(0.37)

    for i, (name, result, person, accent) in enumerate(left):
        y = Inches(1.92) + i * row_gap
        if i % 2 == 0:
            rect(sl, Inches(0.6), y - Inches(0.02), Inches(5.6), row_h, ZEBRA)
        dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.07), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        tx(sl, Inches(0.95), y, Inches(1.6), row_h,
           name, sz=11, c=BLACK, bold=True)
        tx(sl, Inches(2.7), y, Inches(1.7), row_h,
           result, sz=10, c=accent, bold=True)
        tx(sl, Inches(4.5), y, Inches(1.8), row_h,
           person, sz=9, c=MGRAY)

    for i, (name, result, person, accent) in enumerate(right):
        y = Inches(1.92) + i * row_gap
        if i % 2 == 0:
            rect(sl, Inches(6.6), y - Inches(0.02), Inches(6.0), row_h, ZEBRA)
        dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.7), y + Inches(0.07), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        tx(sl, Inches(6.95), y, Inches(1.6), row_h,
           name, sz=11, c=BLACK, bold=True)
        tx(sl, Inches(8.7), y, Inches(1.7), row_h,
           result, sz=10, c=accent, bold=True)
        tx(sl, Inches(10.5), y, Inches(1.8), row_h,
           person, sz=9, c=MGRAY)

    # === LOGO BAR ===
    rect(sl, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.0), NAVY)
    tx(sl, Inches(0.8), Inches(5.95), Inches(5), Inches(0.3),
       "ALSO FEATURED", sz=11, c=RGBColor(0xFF, 0xCC, 0x80), bold=True)

    logos = ["HubSpot", "Notion", "Ramp", "Grafana", "Cursor", "Okta", "Klaviyo", "Uber", "Canva", "Google"]
    pill_colors = [BLUE, PURPLE, TEAL, ORANGE, FOREST, CRIMSON, BLUE, PURPLE, TEAL, ORANGE]
    for i, (name, pc) in enumerate(zip(logos, pill_colors)):
        x = Inches(0.65) + Inches(i * 1.22)
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.3), Inches(1.08), Inches(0.32))
        sh.fill.solid(); sh.fill.fore_color.rgb = pc; sh.line.fill.background()
        tx(sl, x, Inches(6.31), Inches(1.08), Inches(0.3),
           name, sz=9, c=WHITE, bold=True, al=PP_ALIGN.CENTER)

    tx(sl, Inches(0.8), Inches(6.72), Inches(12), Inches(0.2),
       "Source: clay.com  |  SCOUT Framework  |  ISV Intel  |  April 2026",
       sz=8, c=RGBColor(0x90, 0xA0, 0xB8))


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slide_1(prs)
    slide_2(prs)
    out = os.path.join(os.path.expanduser("~"), "Desktop", "clay-intel-v2.pptx")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()

